"""
Vigzone AI - File Processing
=============================
Converts any uploaded file into a form the chat model can consume.

Supported categories
--------------------
Images         → resized / base64-encoded data URI  (PNG, JPG, WEBP, GIF, BMP, TIFF, ICO)
Documents      → extracted text                      (PDF, DOCX, RTF)
Spreadsheets   → table text                          (XLSX, XLSM, CSV, TSV)
Presentations  → slide-by-slide text                 (PPTX)
Data files     → pretty-printed content              (JSON, JSONL, XML, YAML, TOML)
Code / scripts → syntax-highlighted plain text       (py, js, ts, java, c, cpp, cs, go, rs, …)
Archives       → file manifest only                  (ZIP, TAR, TGZ)
Audio / Video  → metadata only, clearly labelled     (MP3, WAV, MP4, MOV, …)
Plain text     → UTF-8 decoded                       (TXT, MD, LOG, INI, ENV, …)
Unknown binary formats are rejected instead of being misrepresented as text.

All extractors share the same _truncate() ceiling so the model's
context window is never blown out.
"""

from __future__ import annotations

import base64
import io
import json
import logging
import zipfile
import tarfile
import os
import struct
from pathlib import Path

from defusedxml import ElementTree as ET
try:
    import magic as _magic
    _HAS_MAGIC = True
except ImportError:
    _HAS_MAGIC = False

from PIL import Image
from pypdf import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation

logger = logging.getLogger("vigzone.files")

try:
    import pypdfium2 as pdfium
    import pytesseract

    _HAS_PDF_OCR = True
except ImportError:
    _HAS_PDF_OCR = False

# ── Constants ──────────────────────────────────────────────────────────────────
MAX_IMAGE_DIMENSION = 1280
IMAGE_JPEG_QUALITY  = 82
MAX_DOC_CHARS       = 20_000   # raised from 15k to support larger files
MAX_OFFICE_UNCOMPRESSED_BYTES = 120 * 1024 * 1024
MAX_OFFICE_ENTRIES = 10_000
MAX_PDF_OCR_PAGES = int(os.getenv("MAX_PDF_OCR_PAGES", "12"))
MAX_PDF_TEXT_PAGES = max(1, min(int(os.getenv("MAX_PDF_TEXT_PAGES", "200")), 1000))
MAX_WORKBOOK_SHEETS = 50
MAX_SHEET_ROWS = 500
MAX_ROW_CELLS = 200
MAX_PRESENTATION_SLIDES = 300
Image.MAX_IMAGE_PIXELS = 40_000_000

# ── Helpers ────────────────────────────────────────────────────────────────────

class FileProcessingError(Exception):
    """Raised when a file can't be read or converted into a usable form."""


def _truncate(text: str) -> tuple[str, bool]:
    text = text.strip()
    if len(text) > MAX_DOC_CHARS:
        return text[:MAX_DOC_CHARS], True
    return text, False


def _sniff_mime(data: bytes) -> str:
    """Return MIME type string. Falls back to 'application/octet-stream'."""
    if _HAS_MAGIC:
        try:
            return _magic.from_buffer(data, mime=True) or "application/octet-stream"
        except Exception:
            pass
    return "application/octet-stream"


def _validate_zip_container(data: bytes, label: str) -> None:
    """Reject malformed or explosively compressed Office/ZIP containers."""

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_OFFICE_ENTRIES:
                raise FileProcessingError(f"{label} contains too many internal files.")
            total_uncompressed = sum(max(0, item.file_size) for item in entries)
            total_compressed = sum(max(1, item.compress_size) for item in entries)
            if total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES:
                raise FileProcessingError(f"{label} expands beyond the safe processing limit.")
            if total_uncompressed > 10 * 1024 * 1024 and total_uncompressed / total_compressed > 250:
                raise FileProcessingError(f"{label} has an unsafe compression ratio.")
    except FileProcessingError:
        raise
    except (zipfile.BadZipFile, OSError) as exc:
        raise FileProcessingError(f"Couldn't open that {label}.") from exc


# ── Image ──────────────────────────────────────────────────────────────────────

def process_image(data: bytes) -> tuple[str, str]:
    """Resize/compress an image and return (data_uri, mime_type)."""
    try:
        img = Image.open(io.BytesIO(data))
        width, height = img.size
        if width <= 0 or height <= 0 or width * height > Image.MAX_IMAGE_PIXELS:
            raise FileProcessingError("That image's pixel dimensions exceed the safe processing limit.")
        img.load()
    except FileProcessingError:
        raise
    except Exception as e:
        raise FileProcessingError("That doesn't look like a readable image file.") from e

    is_png = img.format == "PNG"
    if not is_png and img.mode in ("RGBA", "P", "LA"):
        img = img.convert("RGB")

    if max(width, height) > MAX_IMAGE_DIMENSION:
        scale = MAX_IMAGE_DIMENSION / max(width, height)
        img = img.resize(
            (max(1, int(width * scale)), max(1, int(height * scale))),
            Image.LANCZOS,
        )

    buf = io.BytesIO()
    if is_png:
        img.save(buf, format="PNG", optimize=True)
        mime_type = "image/png"
    else:
        img.save(buf, format="JPEG", quality=IMAGE_JPEG_QUALITY, optimize=True)
        mime_type = "image/jpeg"

    encoded = base64.b64encode(buf.getvalue()).decode("ascii")
    return f"data:{mime_type};base64,{encoded}", mime_type


# ── PDF ────────────────────────────────────────────────────────────────────────

def _ocr_pdf_text(data: bytes, page_count: int) -> str:
    if not _HAS_PDF_OCR or page_count <= 0:
        return ""
    language = os.getenv("OCR_LANGUAGES", "eng").strip() or "eng"
    parts: list[str] = []
    try:
        document = pdfium.PdfDocument(data)
        try:
            for page_index in range(min(len(document), MAX_PDF_OCR_PAGES)):
                page = document[page_index]
                try:
                    bitmap = page.render(scale=1.6)
                    try:
                        image = bitmap.to_pil()
                        text = pytesseract.image_to_string(image, lang=language).strip()
                        if text:
                            parts.append(f"[Page {page_index + 1} — OCR]\n{text}")
                    finally:
                        bitmap.close()
                finally:
                    page.close()
        finally:
            document.close()
    except Exception:
        return ""
    return "\n\n".join(parts)


def extract_pdf_text(data: bytes) -> tuple[str, bool]:
    """Extract text from a PDF. Returns (text, was_truncated).

    Important UX rule:
    Do not reject a valid PDF just because it has no selectable text. Many PDFs
    users upload are image/scanned PDFs (labels, posters, screenshots exported as
    PDF). In that case we accept the file and attach a clear note to the chat
    instead of turning the upload chip red.
    """
    try:
        reader = PdfReader(io.BytesIO(data))
        if getattr(reader, "is_encrypted", False):
            try:
                unlocked = reader.decrypt("")
            except Exception as exc:
                raise FileProcessingError("This PDF is password-protected.") from exc
            if not unlocked:
                raise FileProcessingError("This PDF is password-protected.")
    except FileProcessingError:
        raise
    except Exception as e:
        raise FileProcessingError("Couldn't open that PDF — it may be corrupted or password-protected.") from e

    page_count = 0
    try:
        page_count = len(reader.pages)
    except Exception:
        page_count = 0

    pages = []
    page_limited = page_count > MAX_PDF_TEXT_PAGES
    for page_index in range(min(page_count, MAX_PDF_TEXT_PAGES)):
        try:
            page = reader.pages[page_index]
            t = page.extract_text() or ""
            if t.strip():
                pages.append(f"[Page {page_index + 1}]\n{t}")
        except Exception:
            # Keep going; one broken page should not reject the whole PDF.
            pass

    text = "\n\n".join(pages)
    if not text.strip():
        text = _ocr_pdf_text(data, page_count)
    if not text.strip():
        page_word = "page" if page_count == 1 else "pages"
        limit_note = (
            f" Text extraction was limited to the first {MAX_PDF_TEXT_PAGES} pages."
            if page_limited
            else ""
        )
        fallback = (
            f"[PDF attached: {page_count or 'unknown'} {page_word}]\n"
            "No selectable text or OCR text could be extracted. Vigzone cannot "
            "truthfully analyze the visual contents of this PDF from this attachment. "
            "Upload the relevant pages as PNG/JPEG images for vision analysis."
            + limit_note
        )
        fallback_text, char_limited = _truncate(fallback)
        return fallback_text, bool(page_limited or char_limited)

    if page_limited:
        text += (
            f"\n\n[PDF text extraction limited to the first {MAX_PDF_TEXT_PAGES} "
            f"of {page_count} pages.]"
        )

    truncated_text, char_limited = _truncate(text)
    return truncated_text, bool(page_limited or char_limited)


# ── DOCX ───────────────────────────────────────────────────────────────────────

def extract_docx_text(data: bytes) -> tuple[str, bool]:
    """Extract text from a Word document (.docx)."""
    _validate_zip_container(data, "Word document")
    try:
        doc = Document(io.BytesIO(data))
    except Exception as e:
        raise FileProcessingError("Couldn't open that Word document.") from e

    parts: list[str] = []
    limited = False
    for paragraph in doc.paragraphs[:2000]:
        if paragraph.text.strip():
            parts.append(paragraph.text)
        if sum(len(part) for part in parts) > MAX_DOC_CHARS * 2:
            limited = True
            break
    if len(doc.paragraphs) > 2000:
        limited = True
    rows_seen = 0
    for table in doc.tables[:100]:
        for row in table.rows:
            if rows_seen >= 1000:
                limited = True
                break
            cells = [c.text.strip() for c in row.cells[:MAX_ROW_CELLS]]
            if any(cells):
                parts.append(" | ".join(cells))
            rows_seen += 1
        if limited or sum(len(part) for part in parts) > MAX_DOC_CHARS * 2:
            limited = True
            break
    if len(doc.tables) > 100:
        limited = True

    text = "\n".join(parts)
    if not text.strip():
        raise FileProcessingError("That document appears to be empty.")
    truncated_text, char_limited = _truncate(text)
    return truncated_text, bool(limited or char_limited)


# ── XLSX ───────────────────────────────────────────────────────────────────────

def extract_xlsx_text(data: bytes) -> tuple[str, bool]:
    """Extract a readable table from an Excel workbook (.xlsx / .xlsm)."""
    _validate_zip_container(data, "Excel workbook")
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:
        raise FileProcessingError("Couldn't open that Excel file.") from e

    sections: list[str] = []
    limited = len(wb.sheetnames) > MAX_WORKBOOK_SHEETS
    for sheet_name in wb.sheetnames[:MAX_WORKBOOK_SHEETS]:
        ws = wb[sheet_name]
        rows: list[str] = []
        row_count = 0
        for row in ws.iter_rows(values_only=True):
            if row_count >= MAX_SHEET_ROWS:
                rows.append("… (sheet truncated)")
                limited = True
                break
            if len(row) > MAX_ROW_CELLS:
                limited = True
            cells = [str(c) if c is not None else "" for c in row[:MAX_ROW_CELLS]]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
                row_count += 1
        if rows:
            sections.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        if sum(len(section) for section in sections) > MAX_DOC_CHARS * 2:
            limited = True
            break

    text = "\n\n".join(sections)
    if not text.strip():
        raise FileProcessingError("That spreadsheet appears to be empty.")
    truncated_text, char_limited = _truncate(text)
    return truncated_text, bool(limited or char_limited)


# ── PPTX ───────────────────────────────────────────────────────────────────────

def extract_pptx_text(data: bytes) -> tuple[str, bool]:
    """Extract slide-by-slide text from a PowerPoint file (.pptx)."""
    _validate_zip_container(data, "PowerPoint presentation")
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise FileProcessingError("Couldn't open that PowerPoint file.") from e

    slides: list[str] = []
    limited = len(prs.slides) > MAX_PRESENTATION_SLIDES
    for i, slide in enumerate(prs.slides, 1):
        if i > MAX_PRESENTATION_SLIDES:
            limited = True
            break
        parts: list[str] = []
        if len(slide.shapes) > 500:
            limited = True
        for shape_index, shape in enumerate(slide.shapes):
            if shape_index >= 500:
                break
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
        if sum(len(item) for item in slides) > MAX_DOC_CHARS * 2:
            limited = True
            break

    text = "\n\n".join(slides)
    if not text.strip():
        raise FileProcessingError("That presentation appears to have no readable text.")
    truncated_text, char_limited = _truncate(text)
    return truncated_text, bool(limited or char_limited)


# ── CSV / TSV ──────────────────────────────────────────────────────────────────

def extract_csv_text(data: bytes, delimiter: str = ",") -> tuple[str, bool]:
    """Decode and return CSV/TSV content."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("utf-8", errors="replace")
    lines = text.splitlines()
    # Keep up to 300 rows so we don't blow out context
    if len(lines) > 300:
        preview = "\n".join(lines[:300]) + f"\n… ({len(lines) - 300} more rows)"
        return preview, True
    return _truncate(text)


# ── JSON / JSONL ───────────────────────────────────────────────────────────────

def extract_json_text(data: bytes) -> tuple[str, bool]:
    """Pretty-print JSON or summarise JSONL."""
    try:
        text_raw = data.decode("utf-8", errors="replace")
    except Exception as e:
        raise FileProcessingError("Couldn't decode JSON file.") from e
    try:
        parsed = json.loads(text_raw)
        pretty = json.dumps(parsed, indent=2, ensure_ascii=False)
        return _truncate(pretty)
    except json.JSONDecodeError:
        # Might be JSONL
        lines = text_raw.strip().splitlines()
        results = []
        for line in lines[:50]:
            try:
                results.append(json.dumps(json.loads(line), indent=2, ensure_ascii=False))
            except json.JSONDecodeError:
                results.append(line)
        summary = f"JSONL file ({len(lines)} records). First 50 records:\n\n" + "\n---\n".join(results)
        return _truncate(summary)


# ── XML ────────────────────────────────────────────────────────────────────────

def extract_xml_text(data: bytes) -> tuple[str, bool]:
    """Return XML as-is (text decode), with a structural summary header."""
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception as e:
        raise FileProcessingError("Couldn't decode XML file.") from e
    try:
        root = ET.fromstring(text)
        summary = f"Root element: <{root.tag}>  Children: {len(list(root))}\n\n"
    except ET.ParseError:
        summary = "(XML parse warning — showing raw content)\n\n"
    return _truncate(summary + text)


# ── YAML / TOML ────────────────────────────────────────────────────────────────

def extract_yaml_toml_text(data: bytes) -> tuple[str, bool]:
    """Return YAML or TOML as plain text."""
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception as e:
        raise FileProcessingError("Couldn't decode that config file.") from e
    return _truncate(text)


# ── Archives ───────────────────────────────────────────────────────────────────

def extract_archive_manifest(data: bytes, filename: str) -> tuple[str, bool]:
    """Return a bounded file listing for ZIP/TAR archives (never extract)."""
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".zip" or zipfile.is_zipfile(io.BytesIO(data)):
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                entries = zf.infolist()
                total = len(entries)
                if total > MAX_OFFICE_ENTRIES:
                    raise FileProcessingError("Archive contains too many entries.")
                total_uncompressed = sum(max(0, item.file_size) for item in entries)
                preview = entries[:200]
                listing = "\n".join(
                    f"{item.filename} ({item.file_size} bytes)"
                    for item in preview
                )
                summary = (
                    f"ZIP archive manifest only: {total} file(s), "
                    f"{total_uncompressed} uncompressed bytes declared.\n"
                    "File contents were not opened or analyzed.\n"
                    + ("(showing first 200)\n\n" if total > 200 else "\n")
                    + listing
                )
                return _truncate(summary)
    except FileProcessingError:
        raise
    except (zipfile.BadZipFile, OSError):
        pass

    try:
        with tarfile.open(fileobj=io.BytesIO(data), mode="r:*") as tf:
            total = 0
            total_uncompressed = 0
            preview: list[str] = []
            for member in tf:
                total += 1
                if total > MAX_OFFICE_ENTRIES:
                    raise FileProcessingError("Archive contains too many entries.")
                total_uncompressed += max(0, member.size)
                if len(preview) < 200:
                    preview.append(f"{member.name} ({member.size} bytes)")
            listing = "\n".join(preview)
            summary = (
                f"TAR archive manifest only: {total} file(s), "
                f"{total_uncompressed} uncompressed bytes declared.\n"
                "File contents were not opened or analyzed.\n"
                + ("(showing first 200)\n\n" if total > 200 else "\n")
                + listing
            )
            return _truncate(summary)
    except FileProcessingError:
        raise
    except (tarfile.TarError, OSError):
        pass

    raise FileProcessingError(
        "Couldn't read that archive. Supported manifest formats are ZIP, TAR, and TGZ."
    )


# ── Audio / Video ──────────────────────────────────────────────────────────────

def _read_mp3_id3(data: bytes) -> dict:
    """Read basic ID3v2 tags from MP3 bytes."""
    tags: dict[str, str] = {}
    if data[:3] != b"ID3":
        return tags
    try:
        # ID3v2 frame reading (simplified)
        offset = 10
        while offset < min(len(data), 8192):
            frame_id = data[offset:offset + 4].decode("latin-1", errors="replace")
            if not frame_id.strip() or frame_id[0] not in "ABCDEFGHIJKLMNOPQRSTUVWXYZ0":
                break
            size = struct.unpack(">I", data[offset + 4:offset + 8])[0]
            content = data[offset + 10:offset + 10 + size]
            try:
                tags[frame_id] = content.decode("utf-8", errors="replace").strip("\x00")
            except Exception:
                pass
            offset += 10 + size
    except Exception:
        pass
    return tags


def extract_audio_video_info(data: bytes, filename: str) -> tuple[str, bool]:
    """Return metadata summary for audio/video files."""
    ext = Path(filename).suffix.lower()
    size_kb = len(data) / 1024
    lines = [
        "[Metadata only — media content was not transcribed or visually analyzed]",
        f"File: {filename}",
        f"Size: {size_kb:.1f} KB",
    ]

    if ext == ".mp3":
        tags = _read_mp3_id3(data)
        if tags:
            for key, label in [("TIT2", "Title"), ("TPE1", "Artist"),
                                ("TALB", "Album"), ("TDRC", "Year"),
                                ("TCON", "Genre"), ("TRCK", "Track")]:
                if key in tags and tags[key]:
                    lines.append(f"{label}: {tags[key]}")
        else:
            lines.append("No ID3 tags found.")
    elif ext in {".wav"}:
        if data[:4] == b"RIFF" and data[8:12] == b"WAVE":
            try:
                num_channels = struct.unpack_from("<H", data, 22)[0]
                sample_rate  = struct.unpack_from("<I", data, 24)[0]
                bits         = struct.unpack_from("<H", data, 34)[0]
                lines += [
                    f"Channels: {num_channels}",
                    f"Sample Rate: {sample_rate} Hz",
                    f"Bit Depth: {bits}-bit",
                ]
            except Exception:
                lines.append("WAV header unreadable.")
    else:
        mime = _sniff_mime(data)
        lines.append(f"MIME type: {mime}")
        lines.append("(Full metadata extraction for this format requires ffprobe.)")

    summary = "\n".join(lines)
    return _truncate(summary)


# ── Plain text & code ──────────────────────────────────────────────────────────

# Extensions considered "code"
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".jsx", ".tsx", ".java", ".c", ".cpp", ".cc",
    ".h", ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt",
    ".scala", ".sh", ".bash", ".zsh", ".ps1", ".bat", ".sql", ".r",
    ".m", ".lua", ".dart", ".hs", ".ex", ".exs", ".erl", ".ml", ".fs",
    ".fsx", ".clj", ".cljs", ".groovy", ".pl", ".pm",
}

def extract_plain_text(data: bytes) -> tuple[str, bool]:
    """Decode a plain-text / code file."""
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("utf-8", errors="replace")
    if not text.strip():
        raise FileProcessingError("That file appears to be empty.")
    return _truncate(text)


# ── RTF ────────────────────────────────────────────────────────────────────────

def extract_rtf_text(data: bytes) -> tuple[str, bool]:
    """Very basic RTF → plaintext (strips control words)."""
    import re
    try:
        raw = data.decode("latin-1", errors="replace")
    except Exception as e:
        raise FileProcessingError("Couldn't decode RTF file.") from e
    # Remove RTF control groups and words
    text = re.sub(r"\{[^{}]*\}", "", raw)          # braced groups
    text = re.sub(r"\\[a-z]+\d*\s?", "", text)      # control words
    text = re.sub(r"\\[^a-z]", "", text)            # control symbols
    text = text.replace("\r\n", "\n").strip()
    if not text:
        raise FileProcessingError("No readable text extracted from that RTF file.")
    return _truncate(text)


# ── Universal dispatcher ───────────────────────────────────────────────────────

# Extensions → handler tag
_EXT_MAP: dict[str, str] = {
    # Images
    **{e: "image" for e in (".png", ".jpg", ".jpeg", ".webp", ".gif",
                             ".bmp", ".tiff", ".tif", ".ico")},
    # Documents
    ".pdf":  "pdf",
    ".docx": "docx",
    ".rtf":  "rtf",
    # Spreadsheets
    ".xlsx": "xlsx", ".xlsm": "xlsx",
    ".csv":  "csv",
    ".tsv":  "tsv",
    # Presentations
    ".pptx": "pptx",
    # Data
    ".json": "json", ".jsonl": "json",
    ".xml":  "xml",
    ".yaml": "yaml", ".yml": "yaml",
    ".toml": "yaml",   # same plain-text handler
    # Archives
    ".zip": "archive", ".tar": "archive", ".tgz": "archive",
    # Audio / Video
    **{e: "av" for e in (".mp3", ".wav", ".ogg", ".flac", ".aac",
                          ".mp4", ".mov", ".avi", ".mkv", ".webm", ".m4v", ".m4a")},
    # Plain text & code
    **{e: "text" for e in (".txt", ".md", ".log", ".ini", ".env",
                             ".cfg", ".conf", ".properties",
                             *CODE_EXTENSIONS)},
}

_UNSUPPORTED_HINTS = {
    ".svg": "SVG is not rasterized in this build. Export it as PNG or JPEG.",
    ".doc": "Legacy .doc is not supported. Save it as .docx first.",
    ".odt": "ODT is not supported. Export it as .docx or PDF first.",
    ".xls": "Legacy .xls is not supported. Save it as .xlsx or CSV first.",
    ".ods": "ODS is not supported. Export it as .xlsx or CSV first.",
    ".ppt": "Legacy .ppt is not supported. Save it as .pptx or PDF first.",
    ".7z": "7z archives are not supported. Use ZIP, TAR, or TGZ.",
    ".rar": "RAR archives are not supported. Use ZIP, TAR, or TGZ.",
}


def _looks_like_text(data: bytes) -> bool:
    if not data:
        return False
    sample = data[:8192]
    if b"\x00" in sample:
        return False
    try:
        decoded = sample.decode("utf-8")
    except UnicodeDecodeError:
        return False
    printable = sum(character.isprintable() or character in "\r\n\t" for character in decoded)
    return bool(decoded) and printable / len(decoded) >= 0.9


def process_file(data: bytes, filename: str) -> dict:
    """
    Universal entry point.

    Returns a dict with at minimum:
      kind: "image" | "document" | "archive" | "audio_video" | "unsupported"
      name: filename
    Plus kind-specific fields (data_uri/mime for images; text/truncated for docs).
    """
    ext = Path(filename).suffix.lower()
    if filename.lower().endswith(".tar.gz"):
        ext = ".tgz"
    if ext in _UNSUPPORTED_HINTS:
        raise FileProcessingError(_UNSUPPORTED_HINTS[ext])
    handler = _EXT_MAP.get(ext)

    # MIME sniff fallback when extension is unknown / missing
    if not handler:
        mime = _sniff_mime(data)
        if mime in {"image/png", "image/jpeg", "image/webp", "image/gif", "image/bmp", "image/tiff"}:
            handler = "image"
        elif mime in ("application/pdf",):
            handler = "pdf"
        elif mime in ("application/json", "text/json"):
            handler = "json"
        elif mime.startswith("text/"):
            handler = "text"
        elif mime in ("application/zip",):
            handler = "archive"
        elif _looks_like_text(data):
            handler = "text"
        else:
            raise FileProcessingError(
                "This binary file type is not supported. Convert it to a listed document, image, text, or archive format."
            )

    base = {"name": filename}

    try:
        if handler == "image":
            data_uri, mime = process_image(data)
            result = {**base, "kind": "image", "mime": mime, "data_uri": data_uri}
            if ext == ".gif":
                result.update({
                    "capability": "first_frame_only",
                    "analysis_limited": True,
                    "limitation": "Animated GIFs are analyzed as a single first-frame image.",
                })
            return result

        if handler == "pdf":
            text, trunc = extract_pdf_text(data)
        elif handler == "docx":
            text, trunc = extract_docx_text(data)
        elif handler == "rtf":
            text, trunc = extract_rtf_text(data)
        elif handler == "xlsx":
            text, trunc = extract_xlsx_text(data)
        elif handler == "csv":
            text, trunc = extract_csv_text(data, ",")
        elif handler == "tsv":
            text, trunc = extract_csv_text(data, "\t")
        elif handler == "pptx":
            text, trunc = extract_pptx_text(data)
        elif handler == "json":
            text, trunc = extract_json_text(data)
        elif handler == "xml":
            text, trunc = extract_xml_text(data)
        elif handler == "yaml":
            text, trunc = extract_yaml_toml_text(data)
        elif handler == "archive":
            text, trunc = extract_archive_manifest(data, filename)
            return {
                **base,
                "kind": "archive",
                "text": text,
                "truncated": trunc,
                "capability": "manifest_only",
                "analysis_limited": True,
            }
        elif handler == "av":
            text, trunc = extract_audio_video_info(data, filename)
            return {
                **base,
                "kind": "audio_video",
                "text": text,
                "truncated": trunc,
                "capability": "metadata_only",
                "analysis_limited": True,
            }
        else:
            text, trunc = extract_plain_text(data)

        limited = handler == "pdf" and text.startswith("[PDF attached:")
        return {
            **base,
            "kind": "document",
            "text": text,
            "truncated": trunc,
            "analysis_limited": limited,
        }

    except FileProcessingError:
        raise
    except Exception as exc:
        logger.exception("Unexpected file-processing failure for %s", filename)
        raise FileProcessingError("The file could not be processed safely.") from exc
